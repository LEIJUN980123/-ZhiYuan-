# multi_format_to_json_structured.py
import os
import json
import hashlib
from pathlib import Path
from datetime import datetime
import logging
import io
import zipfile
import sys

# === 第三方库 ===
try:
    from PyPDF2 import PdfReader
except ImportError:
    raise ImportError("请安装: pip install PyPDF2")

try:
    from docx import Document
except ImportError:
    raise ImportError("请安装: pip install python-docx")

try:
    from pptx import Presentation
except ImportError:
    raise ImportError("请安装: pip install python-pptx")

try:
    import pandas as pd
except ImportError:
    raise ImportError("请安装: pip install pandas openpyxl")

try:
    from PIL import Image
    import cv2
    import numpy as np
    import easyocr
except ImportError:
    raise ImportError("请安装: pip install pillow opencv-python-headless numpy easyocr")

import concurrent.futures
from tqdm import tqdm  # 👈 新增：进度条支持

logging.basicConfig(level=logging.INFO, format='%(levelname)s:%(name)s:%(message)s')
logger = logging.getLogger(__name__)

# ========================
# 全局 OCR Reader（每个子进程独立初始化）
# ========================
def get_ocr_reader():
    try:
        reader = easyocr.Reader(['ch_sim', 'en'], gpu=True)
        logger.info("✅ EasyOCR 初始化成功 (GPU 启用)")
        return reader
    except Exception as e:
        logger.warning(f"⚠️ GPU 不可用，回退到 CPU: {e}")
        reader = easyocr.Reader(['ch_sim', 'en'], gpu=False)
        return reader

# ========================
# 工具函数
# ========================

def is_valid_zip(file_path: str) -> bool:
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            zf.testzip()
        return True
    except (zipfile.BadZipFile, OSError, RuntimeError):
        return False

def get_file_fingerprint(file_path: Path) -> str:
    stat = file_path.stat()
    fingerprint_str = f"{file_path.resolve()}|{stat.st_mtime}|{stat.st_size}"
    return hashlib.md5(fingerprint_str.encode()).hexdigest()

def load_processed_index(index_path: Path):
    if index_path.exists():
        try:
            with open(index_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception as e:
            logger.warning(f"索引文件损坏，将重建: {e}")
    return {}

def save_processed_index(index_path: Path, index_data: dict):
    index_path.parent.mkdir(parents=True, exist_ok=True)
    with open(index_path, 'w', encoding='utf-8') as f:
        json.dump(index_data, f, ensure_ascii=False, indent=2)

def image_to_hash(pil_image: Image.Image) -> str:
    if pil_image.mode != 'RGB':
        pil_image = pil_image.convert('RGB')
    img_bytes = io.BytesIO()
    pil_image.save(img_bytes, format='PNG')
    return hashlib.sha256(img_bytes.getvalue()).hexdigest()

def ocr_with_cache(pil_image: Image.Image, image_cache: dict, reader) -> str:
    img_hash = image_to_hash(pil_image)
    if img_hash in image_cache:
        return image_cache[img_hash]
    
    open_cv_image = np.array(pil_image)
    if open_cv_image.shape[-1] == 4:
        open_cv_image = cv2.cvtColor(open_cv_image, cv2.COLOR_RGBA2RGB)
    elif len(open_cv_image.shape) == 2:
        open_cv_image = cv2.cvtColor(open_cv_image, cv2.COLOR_GRAY2RGB)
    elif open_cv_image.shape[-1] == 3:
        open_cv_image = cv2.cvtColor(open_cv_image, cv2.COLOR_RGB2BGR)
    
    try:
        results = reader.readtext(open_cv_image, detail=0, paragraph=True)
        text = '\n'.join(results).strip()
    except Exception as e:
        logger.warning(f"EasyOCR 识别失败: {e}")
        text = ""
    
    image_cache[img_hash] = text
    return text

# ========================
# 文档解析函数
# ========================

def extract_pdf(file_path: str, image_cache: dict, reader) -> list:
    try:
        reader_pdf = PdfReader(file_path)
        sections = []
        for page_num, page in enumerate(reader_pdf.pages):
            extracted = page.extract_text()
            if extracted:
                sections.append({
                    "type": "paragraph",
                    "text": f"Page {page_num + 1}:\n{extracted}"
                })
        return sections
    except Exception as e:
        logger.error(f"PDF 解析失败 ({file_path}): {e}")
        return []

def extract_docx(file_path: str, image_cache: dict, reader) -> list:
    if not is_valid_zip(file_path):
        logger.error(f"无效 DOCX 文件（非 ZIP 格式）: {file_path}")
        return []
    try:
        doc = Document(file_path)
        sections = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            is_heading = (
                len(text) < 60 and 
                any(run.bold for run in para.runs if run.bold is not None)
            )
            if is_heading:
                sections.append({"type": "heading", "level": 2, "text": text})
            else:
                sections.append({"type": "paragraph", "text": text})
        
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    rows.append(cells)
            if rows:
                sections.append({"type": "table", "rows": rows})
        
        image_count = 0
        for rel in doc.part.rels.values():
            if getattr(rel, 'is_external', False):
                continue
            if "image" in rel.target_ref:
                try:
                    image_data = rel.target_part.blob
                    pil_img = Image.open(io.BytesIO(image_data))
                    ocr_result = ocr_with_cache(pil_img, image_cache, reader)
                    if ocr_result:
                        sections.append({
                            "type": "image_ocr",
                            "caption": f"图片 #{image_count + 1}",
                            "text": ocr_result
                        })
                    image_count += 1
                except Exception as e:
                    logger.warning(f"DOCX 图片 OCR 失败: {e}")
        return sections
    except Exception as e:
        logger.error(f"DOCX 解析失败 ({file_path}): {e}")
        return []

def extract_xlsx(file_path: str, image_cache: dict, reader) -> list:
    if not is_valid_zip(file_path):
        logger.error(f"无效 XLSX 文件（非 ZIP 格式）: {file_path}")
        return []
    try:
        sections = []
        excel_file = pd.ExcelFile(file_path, engine='openpyxl')
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(
                file_path,
                sheet_name=sheet_name,
                dtype=str,
                engine='openpyxl',
                keep_default_na=False
            )
            df = df.fillna("")
            rows = [df.columns.tolist()] + df.values.tolist()
            non_empty_rows = [
                [str(cell).strip() for cell in row] 
                for row in rows 
                if any(str(cell).strip() for cell in row)
            ]
            if non_empty_rows:
                sections.append({
                    "type": "table",
                    "caption": f"Sheet: {sheet_name}",
                    "rows": non_empty_rows
                })
        return sections
    except Exception as e:
        logger.error(f"XLSX 解析失败 ({file_path}): {e}")
        return []

def extract_pptx(file_path: str, image_cache: dict, reader) -> list:
    if not is_valid_zip(file_path):
        logger.error(f"无效 PPTX 文件（非 ZIP 格式）: {file_path}")
        return []
    try:
        prs = Presentation(file_path)
        sections = []
        for slide_num, slide in enumerate(prs.slides):
            slide_sections = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_sections.append({
                        "type": "paragraph",
                        "text": shape.text.strip()
                    })
            
            image_count = 0
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        pil_img = Image.open(io.BytesIO(image_bytes))
                        ocr_result = ocr_with_cache(pil_img, image_cache, reader)
                        if ocr_result:
                            slide_sections.append({
                                "type": "image_ocr",
                                "caption": f"Slide {slide_num + 1} 图片 #{image_count + 1}",
                                "text": ocr_result
                            })
                        image_count += 1
                    except Exception as e:
                        logger.warning(f"PPTX 图片 OCR 失败 (Slide {slide_num + 1}): {e}")
            
            if slide_sections:
                sections.append({
                    "type": "slide",
                    "slide_number": slide_num + 1,
                    "content": slide_sections
                })
        return sections
    except Exception as e:
        logger.error(f"PPTX 解析失败 ({file_path}): {e}")
        return []

# ========================
# 多进程安全的提取函数
# ========================

def extract_file_safe(args):
    file_path, cache_path = args
    reader = get_ocr_reader()
    
    image_cache = {}
    if os.path.exists(cache_path):
        try:
            with open(cache_path, 'r', encoding='utf-8') as f:
                image_cache = json.load(f)
        except Exception as e:
            pass  # 忽略缓存加载失败

    try:
        ext = file_path.lower().split('.')[-1]
        if ext == 'pdf':
            sections = extract_pdf(file_path, image_cache, reader)
            file_type = 'pdf'
        elif ext == 'docx':
            sections = extract_docx(file_path, image_cache, reader)
            file_type = 'docx'
        elif ext == 'xlsx':
            sections = extract_xlsx(file_path, image_cache, reader)
            file_type = 'xlsx'
        elif ext == 'pptx':
            sections = extract_pptx(file_path, image_cache, reader)
            file_type = 'pptx'
        else:
            return None

        if not sections:
            return None

        return {
            "filename": os.path.basename(file_path),
            "file_type": file_type,
            "sections": sections,
            "source_path": os.path.abspath(file_path)
        }
    except Exception as e:
        logger.error(f"处理 {file_path} 时发生未预期错误: {e}")
        return None

# ========================
# 主流程（带进度条）
# ========================

def batch_convert_to_json_incremental(
    input_dir: str,
    output_file: str,
    index_file: str = "output/processed_index.json",
    image_cache_file: str = "output/image_ocr_cache.json",
    max_workers: int = 4
):
    input_path = Path(input_dir)
    output_path = Path(output_file)
    index_path = Path(index_file)
    cache_path = Path(image_cache_file)

    if not input_path.exists():
        raise FileNotFoundError(f"输入目录不存在: {input_dir}")

    existing_results = []
    if output_path.exists():
        try:
            with open(output_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                existing_results = data if isinstance(data, list) else []
        except Exception as e:
            logger.warning(f"输出文件损坏，将重建: {e}")

    file_map = {item["source_path"]: item for item in existing_results if "source_path" in item}
    processed_index = load_processed_index(index_path)

    supported_exts = ('.pdf', '.docx', '.xlsx', '.pptx')
    all_files = [f for f in input_path.rglob('*') if f.is_file() and f.suffix.lower() in supported_exts]

    files_to_process = []
    updated_results = []
    skipped = 0

    for file_path in all_files:
        abs_path = str(file_path.resolve())
        current_fingerprint = get_file_fingerprint(file_path)

        if abs_path in processed_index:
            if processed_index[abs_path]["fingerprint"] == current_fingerprint:
                if abs_path in file_map:
                    updated_results.append(file_map[abs_path])
                    skipped += 1
                    continue

        files_to_process.append(str(file_path))

    new_results = []
    if files_to_process:
        logger.info(f"🔄 将并行处理 {len(files_to_process)} 个文件（进程数: {max_workers}）...")
        
        tasks = [(fp, str(cache_path)) for fp in files_to_process]
        
        with concurrent.futures.ProcessPoolExecutor(max_workers=max_workers) as executor:
            futures = [executor.submit(extract_file_safe, task) for task in tasks]
            
            # 🎯 关键：使用 tqdm 显示实时进度条
            for future in tqdm(
                concurrent.futures.as_completed(futures),
                total=len(futures),
                desc="📄 处理文档",
                unit="文件",
                ncols=80,
                bar_format="{l_bar}{bar}| {n_fmt}/{total_fmt} [{elapsed}<{remaining}, {rate_fmt}]"
            ):
                try:
                    doc = future.result()
                    if doc:
                        new_results.append(doc)
                        abs_path = doc["source_path"]
                        current_fingerprint = get_file_fingerprint(Path(abs_path))
                        processed_index[abs_path] = {
                            "fingerprint": current_fingerprint,
                            "processed_at": datetime.now().isoformat()
                        }
                except Exception as e:
                    logger.error(f"多进程任务异常: {e}")

    updated_results.extend(new_results)
    total = len(updated_results)
    new_or_changed = len(new_results)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(updated_results, f, ensure_ascii=False, indent=2)

    save_processed_index(index_path, processed_index)

    logger.info(f"\n✅ 处理完成！")
    logger.info(f"   - 总文件数: {total}")
    logger.info(f"   - 新增或更新: {new_or_changed}")
    logger.info(f"   - 跳过未变: {skipped}")
    logger.info(f"📁 输出文件: {output_path.absolute()}")
    logger.info(f"📊 索引文件: {index_path.absolute()}")
    logger.info(f"🖼️  图片缓存: {cache_path.absolute()}")

# ========================
# 入口（预下载模型 + 启动）
# ========================

if __name__ == "__main__":
    print("📥 正在预下载 EasyOCR 模型（仅首次运行需要，约1-5分钟）...")
    try:
        _ = easyocr.Reader(['ch_sim', 'en'], gpu=False)
        print("✅ EasyOCR 模型已缓存到本地")
    except Exception as e:
        print(f"⚠️ 模型预下载可能失败（子进程将重试）: {e}")

    INPUT_DIR = "input_docs"
    OUTPUT_FILE = "output/structured_docs.json"
    INDEX_FILE = "output/processed_index.json"
    IMAGE_CACHE_FILE = "output/image_ocr_cache.json"

    workers = min(6, os.cpu_count() or 4)

    batch_convert_to_json_incremental(
        input_dir=INPUT_DIR,
        output_file=OUTPUT_FILE,
        index_file=INDEX_FILE,
        image_cache_file=IMAGE_CACHE_FILE,
        max_workers=workers
    )